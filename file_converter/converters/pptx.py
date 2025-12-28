"""
file_converter/converters/pptx.py
PPTX 文件转换器 - 使用 python-pptx 提取幻灯片内容
"""

import io
import logging
from typing import BinaryIO, List

from pptx import Presentation
from pptx.util import Pt

from file_converter.converters.base import BaseConverter, ConversionResult

logger = logging.getLogger(__name__)


class PptxConverter(BaseConverter):
    """
    PPTX 转 Markdown 转换器
    
    功能：
    - 每页幻灯片生成二级标题 "## Slide N: 标题"
    - 提取文本框内容
    - 备注以引用块格式附加
    - 图片/图表标注为 [图片/图表略]
    """
    
    supported_extensions = ["pptx", "ppt"]
    
    def convert(self, file: BinaryIO, filename: str) -> ConversionResult:
        """将 PPTX 转换为 Markdown"""
        warnings: List[str] = []
        content_parts: List[str] = []
        
        try:
            prs = Presentation(file)
            total_slides = len(prs.slides)
            logger.info(f"Processing PPTX: {filename}, {total_slides} slides")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = self._process_slide(slide, slide_num, warnings)
                content_parts.append(slide_content)
        
        except Exception as e:
            logger.error(f"PPTX conversion failed: {filename}, error: {e}")
            raise ValueError(f"PPTX 转换失败: {str(e)}")
        
        content = "\n\n---\n\n".join(content_parts)
        
        if not content.strip():
            warnings.append("PPTX 未提取到任何内容")
        
        return ConversionResult(content=content, warnings=warnings)
    
    def _process_slide(self, slide, slide_num: int, warnings: List[str]) -> str:
        """处理单个幻灯片"""
        parts: List[str] = []
        
        # 提取标题
        title = self._get_slide_title(slide)
        if title:
            parts.append(f"## Slide {slide_num}: {title}")
        else:
            parts.append(f"## Slide {slide_num}")
        
        # 提取文本内容
        text_content = self._extract_text_content(slide, warnings)
        if text_content:
            parts.append(text_content)
        
        # 提取备注
        notes = self._get_notes(slide)
        if notes:
            parts.append("\n> **备注：**")
            for line in notes.split('\n'):
                if line.strip():
                    parts.append(f"> {line.strip()}")
        
        return "\n\n".join(parts)
    
    def _get_slide_title(self, slide) -> str:
        """获取幻灯片标题"""
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        
        # 尝试从第一个文本框获取标题
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) < 100:
                    return text
        
        return ""
    
    def _extract_text_content(self, slide, warnings: List[str]) -> str:
        """提取幻灯片文本内容"""
        lines: List[str] = []
        has_image = False
        has_chart = False
        
        for shape in slide.shapes:
            # 跳过标题（已单独处理）
            if shape == slide.shapes.title:
                continue
            
            # 检查图片
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                has_image = True
                continue
            
            # 检查图表
            if shape.has_chart:
                has_chart = True
                continue
            
            # 提取文本框内容
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = self._process_paragraph(paragraph)
                    if text:
                        lines.append(text)
            
            # 提取表格
            if shape.has_table:
                table_md = self._table_to_markdown(shape.table)
                if table_md:
                    lines.append(table_md)
        
        # 添加图片/图表标记
        if has_image:
            lines.append("\n[图片略]")
        if has_chart:
            lines.append("\n[图表略]")
        
        return "\n".join(lines)
    
    def _process_paragraph(self, paragraph) -> str:
        """处理段落，转换格式"""
        text = paragraph.text.strip()
        if not text:
            return ""
        
        # 检查列表级别
        level = paragraph.level
        if level > 0:
            indent = "  " * level
            text = f"{indent}- {text}"
        
        # 处理加粗和斜体
        formatted_parts: List[str] = []
        for run in paragraph.runs:
            run_text = run.text
            if run.font.bold:
                run_text = f"**{run_text}**"
            if run.font.italic:
                run_text = f"*{run_text}*"
            formatted_parts.append(run_text)
        
        if formatted_parts:
            return "".join(formatted_parts)
        
        return text
    
    def _get_notes(self, slide) -> str:
        """获取幻灯片备注"""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
        return ""
    
    def _table_to_markdown(self, table) -> str:
        """将 PPTX 表格转换为 Markdown"""
        if not table.rows:
            return ""
        
        lines: List[str] = []
        
        # 获取列数
        col_count = len(table.columns)
        
        for row_idx, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                text = cell.text.strip().replace("|", "\\|").replace("\n", " ")
                cells.append(text)
            
            lines.append("| " + " | ".join(cells) + " |")
            
            # 在第一行后添加分隔行
            if row_idx == 0:
                lines.append("| " + " | ".join(["---"] * col_count) + " |")
        
        return "\n".join(lines)
